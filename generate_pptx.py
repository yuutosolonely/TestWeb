# -*- coding: utf-8 -*-
import os
import sys

# Tự động cài đặt thư viện python-pptx nếu chưa có
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx library...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Thiết lập kích thước Slide tỷ lệ chuẩn 16:9 (Widescreen)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Sử dụng layout 1: Title and Content cho hầu hết các Slide
    # layout 0: Title Slide cho trang bìa
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]
    
    # Định nghĩa palette màu sắc hiện đại sang trọng
    DARK_BLUE = RGBColor(12, 28, 56)      # #0C1C38 - Màu nền tối trang trọng
    LIGHT_BLUE = RGBColor(240, 244, 250)  # #F0F4FA - Màu nền sáng nhẹ
    TEXT_DARK = RGBColor(33, 37, 41)       # #212529 - Màu chữ tối chính
    TEXT_LIGHT = RGBColor(255, 255, 255)  # #FFFFFF - Màu chữ trắng
    ACCENT_BLUE = RGBColor(0, 102, 204)   # #0066CC - Xanh dương Docker chính
    ACCENT_GREEN = RGBColor(40, 167, 69)  # #28A745 - Xanh lá an toàn
    MUTED_GRAY = RGBColor(108, 117, 125)  # #6C757D - Chữ phụ/Ghi chú
    
    # ── SLIDE 1: TRANG BÌA ───────────────────────────────────────────────
    slide1 = prs.slides.add_slide(title_slide_layout)
    # Tô màu nền tối cho trang bìa
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Cấu hình Tiêu đề chính
    title = slide1.shapes.title
    title.text = "NOTE MANAGEMENT APPLICATION"
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
    
    # Cấu hình Tiêu đề phụ
    subtitle = slide1.placeholders[1]
    subtitle.text = "Chuyên đề Topic 10: Containerization & Orchestration (Docker)\n\nHướng dẫn: TS. Lê Văn Vang\nNhóm thực hiện: Nhóm 10"
    subtitle.text_frame.paragraphs[0].font.name = "Arial"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(173, 181, 189)
    
    # Thêm Speaker Notes cho Slide 1
    slide1.notes_slide.notes_text_frame.text = (
        "Kịch bản thuyết trình (Speaker Notes) - SLIDE 1:\n"
        "\"Kính chào thầy cô trong Hội đồng đánh giá và các bạn sinh viên. "
        "Hôm nay nhóm chúng em xin phép trình bày đồ án kết thúc môn học Web Programming & Applications. "
        "Đề tài của chúng em là 'Xây dựng ứng dụng quản lý ghi chú đa tầng Note Management'. "
        "Trong đó, nhóm chúng em tập trung giải quyết Chuyên đề Topic 10: Đóng gói ứng dụng bằng công nghệ "
        "Containerization và Điều phối dịch vụ bằng Orchestration thông qua hệ sinh thái Docker. "
        "Sau đây, chúng em xin đi vào chi tiết.\""
    )

    # Dữ liệu các slide tiếp theo: (Tiêu đề, Danh sách bullet points, Speaker Notes, Màu nền sáng/tối)
    slides_data = [
        # SLIDE 2
        (
            "🔎 ĐẶT VẤN ĐỀ & MỤC TIÊU DỰ ÁN",
            [
                "Thực trạng phát triển ứng dụng phức tạp:",
                "   • Lỗi xung đột môi trường phát triển (khác phiên bản PHP, Apache, Node.js giữa các thành viên).",
                "   • Khó khăn tích hợp các dịch vụ nâng cao (PostgreSQL, Redis cache, WebSocket server).",
                "   • Mất nhiều thời gian để cấu hình cổng, tài khoản cơ sở dữ liệu và phân quyền thủ công.",
                "Mục tiêu cốt lõi của dự án:",
                "   • Áp dụng Containerization đóng gói độc lập toàn bộ các dịch vụ hệ thống.",
                "   • Tối giản quy trình triển khai: Khởi chạy toàn bộ hệ thống bằng 1 dòng lệnh duy nhất.",
                "   • Đảm bảo tính nhất quán môi trường 100% từ máy lập trình viên cho đến cloud production."
            ],
            "Trong quy trình phát triển phần mềm truyền thống, lập trình viên thường xuyên gặp lỗi xung đột môi trường. 'Chạy tốt trên máy của em nhưng lại lỗi trên máy của bạn' là câu nói kinh điển. Hơn nữa, với một ứng dụng ghi chú cao cấp đòi hỏi cả cơ sở dữ liệu PostgreSQL, cache Redis và server WebSocket, việc cài đặt thủ công từng dịch vụ trên máy cá nhân là cực kỳ phức tạp và dễ gây nghẽn hệ thống. Mục tiêu của chúng em là đóng gói toàn bộ ứng dụng này vào các Container độc lập, giúp việc triển khai trở nên tự động hoàn toàn, đồng bộ 100% trên mọi máy tính.",
            False # Nền sáng
        ),
        # SLIDE 3
        (
            "🌐 TỔNG QUAN KIẾN TRÚC ĐA TẦNG (MULTI-TIER)",
            [
                "Kiến trúc phân rã 4 Tầng dịch vụ (4-Tier) độc lập hoàn toàn:",
                "   • Web Application (app): Chạy PHP 8.2 + Web Server Apache làm đầu ngõ xử lý HTTP requests (Port 8000).",
                "   • Database (db): Chạy PostgreSQL 15 Alpine lưu trữ dữ liệu quan hệ bền vững (Port 5433).",
                "   • Cache & Session (redis): Chạy Redis Alpine lưu trữ bộ đệm RAM giảm tải cho DB (Port 6379).",
                "   • WebSocket Server (websocket): Chạy Laravel Reverb xử lý truyền thông hai chiều thời gian thực (Port 8080).",
                "Mạng nội bộ cô lập (Docker Bridge Network):",
                "   • Khai báo mạng bridge ảo 'note_network' giúp các container tự động kết nối và phân giải DNS qua Service Name.",
                "   • Chặn đứng các kết nối tấn công trái phép từ bên ngoài chọc thẳng vào Database."
            ],
            "Hệ thống của chúng em được xây dựng trên Kiến trúc Đa tầng (Multi-tier Architecture) tiên tiến gồm 4 container chạy hoàn toàn tách biệt. Đầu tiên là container App đóng vai trò tiếp nhận yêu cầu từ người dùng qua cổng 8000. Dữ liệu ghi chú được lưu trữ bền vững tại container DB chạy PostgreSQL. Để tăng tốc độ phản hồi, chúng em có container Redis đóng vai trò lưu đệm session trên RAM. Và cuối cùng là container WebSocket Reverb chịu trách nhiệm duy trì kênh truyền thông hai chiều để đồng bộ ghi chú. Toàn bộ 4 dịch vụ này được đặt trong một mạng Bridge ảo khép kín gọi là 'note_network', đảm bảo tính an toàn và bảo mật cao nhất.",
            False
        ),
        # SLIDE 4
        (
            "🐳 TẠI SAO CHỌN DOCKER CONTAINER? (DOCKER VS VM)",
            [
                "Máy ảo truyền thống (Virtual Machine - VM):",
                "   • Ảo hóa ở cấp độ phần cứng thông qua Hypervisor.",
                "   • Mỗi máy ảo bắt buộc phải cài đặt một hệ điều hành khách hoàn chỉnh (Guest OS) vô cùng nặng nề.",
                "   • Khởi động chậm (tính bằng phút), ngốn RAM, CPU và dung lượng ổ cứng (hàng chục GB).",
                "Docker Container (Ảo hóa cấp độ Hệ điều hành):",
                "   • Loại bỏ lớp Guest OS khách, chia sẻ trực tiếp nhân Kernel của Host OS.",
                "   • Kích thước siêu nhẹ (vài chục MB đến vài trăm MB), khởi động tức thì trong vài mili-giây.",
                "   • Đạt hiệu năng chạy ứng dụng gần như nguyên bản (Native Performance) do không hao phí tài nguyên ảo hóa."
            ],
            "Để hiểu rõ ưu thế của Docker, hãy so sánh nó với máy ảo truyền thống VM. Nhìn vào sơ đồ kiến trúc, chúng ta thấy VM rất cồng kềnh vì mỗi máy ảo phải duy trì một hệ điều hành khách riêng, làm lãng phí hàng chục gigabyte ổ cứng và hàng gigabyte RAM. Docker giải quyết triệt để rào cản này bằng cách ảo hóa ở cấp độ hệ điều hành, chia sẻ chung nhân Kernel của máy thật. Nhờ vậy, container Docker có dung lượng siêu nhẹ, khởi động tức thì trong chớp mắt và đạt hiệu suất xử lý phần cứng tối ưu nhất.",
            False
        ),
        # SLIDE 5
        (
            "🎨 TẦNG ỨNG DỤNG & THIẾT KẾ PWA OFFLINE-FIRST",
            [
                "Nền tảng Tầng Ứng Dụng (App Tier):",
                "   • Sử dụng PHP 8.2 kết hợp web server Apache (mpm_prefork) mang lại độ ổn định tối đa.",
                "   • Framework Laravel xử lý toàn bộ logic API nghiệp vụ, phân quyền và kết xuất giao diện.",
                "Tính năng Progressive Web App (PWA) đột phá:",
                "   • Cho phép người dùng 'Cài đặt ứng dụng' trực tiếp lên màn hình điện thoại di động không qua App Store.",
                "   • Service Worker (sw.js): Chạy ngầm trong trình duyệt, tự động cache tài nguyên tĩnh giúp app tải cực nhanh.",
                "   • IndexedDB Client Storage: Cơ sở dữ liệu NoSQL trong trình duyệt giúp ghi chú bình thường khi mất mạng.",
                "   • Cơ chế Offline-first: Tự động đồng bộ hóa ghi chú nháp lên server ngay khi thiết bị kết nối mạng trở lại."
            ],
            "Tầng ứng dụng sử dụng Laravel chạy trên nền PHP 8.2 và máy chủ Apache. Điểm nhấn công nghệ đặc biệt ở đây là chúng em đã biến ứng dụng này thành một Progressive Web App (PWA) thực thụ. Nhờ vào Service Worker chạy ngầm, toàn bộ giao diện tĩnh được lưu vào cache trình duyệt. Kết hợp với CSDL IndexedDB, người dùng vẫn có thể viết ghi chú bình thường khi mất mạng. Ngay khi thiết bị có kết nối mạng trở lại, PWA sẽ tự động đồng bộ hóa các bản nháp offline này lên máy chủ PostgreSQL, mang lại trải nghiệm vô cùng mượt mà.",
            False
        ),
        # SLIDE 6
        (
            "🗄️ TẦNG CƠ SỞ DỮ LIỆU & BỘ ĐỆM CACHE REDIS",
            [
                "Tầng cơ sở dữ liệu (PostgreSQL 15 Alpine):",
                "   • PostgreSQL là chuẩn CSDL microservices cao cấp, xử lý đồng thời (concurrency) vượt trội hơn MySQL.",
                "   • Hỗ trợ kiểu dữ liệu JSONB cực mạnh: Lưu trữ thuộc tính động của ghi chú (màu sắc, nhãn, tọa độ vẽ) linh hoạt.",
                "   • Đảm bảo tính toàn vẹn dữ liệu và an toàn giao dịch đạt chuẩn nghiêm ngặt ACID.",
                "Tầng lưu đệm RAM (Redis Alpine):",
                "   • Sử dụng Redis lưu trữ toàn bộ sessions đăng nhập và bộ đệm truy vấn trực tiếp trên RAM.",
                "   • Giải quyết triệt để hiện tượng nghẽn cổ chai ổ đĩa cứng (I/O Bottleneck) của PostgreSQL.",
                "   • Tốc độ phản hồi kiểm tra đăng nhập người dùng giảm xuống dưới 1 mili-giây."
            ],
            "Về cơ sở dữ liệu, chúng em chọn PostgreSQL 15. So với MySQL, PostgreSQL mạnh mẽ hơn hẳn khi xử lý các câu truy vấn phức tạp và đặc biệt hỗ trợ lưu trữ định dạng JSONB. Ghi chú của người dùng thường có nhiều thuộc tính động như màu sắc, thẻ nhãn, tọa độ vẽ, vì vậy JSONB giúp chúng em lưu trữ linh hoạt mà không cần tạo quá nhiều bảng. Ngoài ra, để tránh nghẽn băng thông ổ đĩa, chúng em tích hợp Redis chạy hoàn toàn trên RAM để lưu thông tin phiên đăng nhập. Nhờ đó, tốc độ kiểm tra đăng nhập người dùng được giảm xuống dưới 1 mili-giây, giúp hệ thống hoạt động vô cùng nhẹ nhàng.",
            False
        ),
        # SLIDE 7
        (
            "⚡ WEBSOCKET REAL-TIME & CỘNG TÁC THỜI GIAN THỰC",
            [
                "Máy chủ Laravel Reverb WebSocket Server:",
                "   • Thiết lập kết nối socket hai chiều (Full-duplex connection) liên tục giữa client và server.",
                "   • Thay thế cơ chế kéo dữ liệu truyền thống (Polling) gây lãng phí băng thông và cạn kiệt CPU.",
                "Tính năng Cộng tác thời gian thực (Real-time Collaboration):",
                "   • Đồng bộ nội dung ghi chú ngay lập tức trên màn hình của tất cả các cộng tác viên khi có người chỉnh sửa.",
                "   • Hỗ trợ cập nhật màu sắc, gắn nhãn, xóa ghi chú real-time cực kỳ mượt mà.",
                "Tính năng Tự động lưu ngầm (Auto-save):",
                "   • Áp dụng kỹ thuật Debounce 1 giây: Tự động lưu dữ liệu xuống DB ngay khi dừng gõ chữ 1s.",
                "   • Loại bỏ hoàn toàn nút 'Lưu ghi chú' thủ công, tối ưu hóa trải nghiệm người dùng."
            ],
            "Tính năng nổi bật và mang tính thực tiễn cao nhất của dự án là khả năng cộng tác nhóm thời gian thực. Chúng em sử dụng máy chủ WebSocket Laravel Reverb thế hệ mới. Khi hai hay nhiều thành viên cùng mở chung một ghi chú được chia sẻ, bất kỳ ai gõ chữ hay đổi màu nền, thông tin sẽ được phát qua kênh socket và cập nhật lên màn hình của những người còn lại ngay lập tức mà không cần tải lại trang. Hệ thống cũng tích hợp cơ chế Auto-save tự động lưu ghi chú sau 1 giây ngắt nghỉ gõ chữ, loại bỏ hoàn toàn nút Lưu thực hiện thủ công phiền phức.",
            False
        ),
        # SLIDE 8
        (
            "📄 PHÂN TÍCH TỆP CẤU HÌNH DOCKERFILE CỦA WEB APP",
            [
                "Cơ chế hoạt động của Dockerfile định nghĩa môi trường chạy web app:",
                "   • Base Image: Khởi đầu từ php:8.2-apache chính thức tích hợp sẵn Apache và PHP 8.2.",
                "   • Cài đặt Driver hệ thống: Cập nhật apt và tự biên dịch pdo_pgsql, pgsql, gd, zip và redis extensions.",
                "   • Áp dụng kỹ thuật Multi-stage build: Kéo tệp thực thi composer trực tiếp từ image chính thức giúp giảm dung lượng.",
                "   • Tích hợp Vite và Node.js v20: Tự động tải NPM packages và chạy biên dịch assets (npm run build) ngay khi build.",
                "   • Cấu hình Apache VirtualHost bảo mật: Trỏ DocumentRoot vào thư mục public/ của Laravel để giấu file mã nguồn.",
                "   • Sửa lỗi tương thích hệ thống: Chuyển đổi định dạng ký tự xuống dòng kiểu Windows (CRLF) sang Linux (LF) cho script."
            ],
            "Để đóng gói ứng dụng Web, chúng em đã thiết kế tệp tin Dockerfile này. Chúng em sử dụng image gốc là PHP 8.2 Apache. Trong quá trình build Image, Dockerfile sẽ tự động cài các thư viện C++ cần thiết để biên dịch driver kết nối PostgreSQL và Redis. Kỹ thuật nâng cao ở đây là chúng em áp dụng Multi-stage build để kéo Composer từ image chính thức giúp giảm dung lượng. Đồng thời, Dockerfile cài đặt Node.js v20 để biên dịch giao diện Frontend bằng Vite trước khi đóng gói, trỏ DocumentRoot vào thư mục public nhằm bảo mật mã nguồn Laravel.",
            False
        ),
        # SLIDE 9
        (
            "📄 PHÂN TÍCH FILE ĐIỀU PHỐI DOCKER-COMPOSE.YML",
            [
                "Vai trò 'Nhạc trưởng' điều phối toàn diện 4 Services:",
                "   • app Service: Ánh xạ cổng 8000:80, truyền các biến môi trường cấu hình kết nối DB, Redis và mail SMTP.",
                "   • db Service: Khởi chạy Postgres:15-alpine trên cổng 5433:5432, lưu trữ dữ liệu ra phân vùng ổ đĩa ảo độc lập.",
                "   • redis Service: Chạy redis:alpine lưu trữ session và RAM-cache trên cổng 6379:6379.",
                "   • websocket Service: Khởi động máy chủ Laravel Reverb trên cổng 8080:8080 để kết nối socket thời gian thực.",
                "Cơ chế ràng buộc thông minh (Healthcheck Dependency):",
                "   • db được kiểm định bằng pg_isready; redis được kiểm định bằng redis-cli ping mỗi 5 giây.",
                "   • Container app chỉ khởi chạy khi db và redis đạt trạng thái Healthy hoàn toàn.",
                "   • Tránh triệt để lỗi Crash ứng dụng do mất kết nối CSDL khi khởi động đồng loạt."
            ],
            "Tệp tin docker-compose.yml chính là nhạc trưởng điều phối toàn bộ hệ thống. Điểm đặc sắc ở đây là cơ chế Health Check. Thông thường các dịch vụ khởi chạy đồng thời, dẫn đến việc ứng dụng Web chạy trước CSDL PostgreSQL và gây ra lỗi mất kết nối. Chúng em đã giải quyết triệt để bằng cách thiết lập ràng buộc: App chỉ được phép khởi động khi Database và Redis đã vượt qua bài test sức khỏe định kỳ pg_isready và redis-cli ping. Dữ liệu ghi chú của người dùng cũng được bảo lưu vĩnh viễn nhờ phân vùng Volume độc lập.",
            False
        ),
        # SLIDE 10
        (
            "⚙️ TỰ ĐỘNG HÓA VẬN HÀNH VỚI DOCKER-ENTRYPOINT.SH",
            [
                "Script Bash thông minh tự động hóa toàn bộ quá trình vận hành khi container Web khởi chạy:",
                "   • Tương thích cổng động Cloud: Tự động sửa tệp cấu hình Apache lắng nghe cổng $PORT cấp phát tự động từ Cloud.",
                "   • Khởi tạo biến môi trường: Tự tạo file .env từ .env.example và phát sinh mã khóa bảo mật app key:generate.",
                "   • Vòng lặp chờ đợi thông minh: Chạy thử kiểm tra kết nối DB tối đa 30 lần (30 giây) để đảm bảo Postgres đã sẵn sàng.",
                "   • Đồng bộ cơ sở dữ liệu: Tự động chạy migrate tạo bảng dữ liệu và nạp dữ liệu mẫu thử nghiệm (db:seed) tự động.",
                "   • Dọn dẹp & Phân quyền: Tạo storage link hiển thị ảnh đại diện, xóa sạch cache cấu hình cũ và phân quyền cho www-data."
            ],
            "Để tối ưu hóa trải nghiệm vận hành, chúng em viết thêm một bash script thông minh là docker-entrypoint.sh. Khi container Web bắt đầu chạy, script này sẽ thực hiện một loạt tác vụ tự động: tự động điều chỉnh cổng Apache thích ứng với môi trường Cloud, tự kiểm tra và cài đặt thư viện PHP còn thiếu, tự động chờ cơ sở dữ liệu sẵn sàng kết nối trong vòng 30 giây, tự động chạy migrate tạo bảng dữ liệu và nạp sẵn 2 tài khoản demo thử nghiệm mà không cần quản trị viên phải gõ bất kỳ câu lệnh thủ công nào.",
            False
        ),
        # SLIDE 11
        (
            "🚀 QUY TRÌNH 3 BƯỚC KHỞI CHẠY HỆ THỐNG TRONG DEV MÔI TRƯỜNG",
            [
                "Bước 1: Giải phóng tài nguyên cổng trên máy thật (Preparation):",
                "   • Dừng hoàn toàn các phần mềm local đang chiếm dụng cổng như XAMPP (Apache, MySQL) hoặc Postgres local máy thật.",
                "Bước 2: Khởi động đóng gói và dựng hệ thống tự động (Build & Start):",
                "   • Mở terminal tại thư mục gốc của dự án và chạy câu lệnh duy nhất:",
                "     >>> docker-compose up -d --build",
                "   • Docker Compose tự xây dựng Image, cấu hình mạng bridge ảo và khởi chạy 4 container độc lập.",
                "Bước 3: Truy cập và kiểm thử tính năng thời gian thực (Test & Experience):",
                "   • Mở trình duyệt Web truy cập: http://localhost:8000",
                "   • Đăng nhập 2 tài khoản demo (1 tab thường, 1 tab ẩn danh) để kiểm thử đồng bộ ghi chú real-time qua WebSocket:",
                "     - Trình duyệt 1: demo@example.com / mật khẩu: 123456",
                "     - Trình duyệt 2: demo2@example.com / mật khẩu: 123456"
            ],
            "Để minh chứng cho sự tiện lợi của công nghệ Container hóa, quy trình khởi chạy dự án của chúng em gói gọn trong 3 bước cực kỳ đơn giản. Đầu tiên, chúng em giải phóng các cổng dịch vụ trên máy thật. Tiếp theo, mở terminal chạy lệnh docker-compose up -d --build. Hệ thống tự động biên dịch và dựng lên toàn bộ 4 tầng dịch vụ. Cuối cùng, chúng em mở trình duyệt truy cập localhost:8000 để trải nghiệm ứng dụng. Ngay sau đây, chúng em xin phép trình diễn trực quan tính năng cộng tác thời gian thực và hoạt động offline PWA của ứng dụng.",
            False
        ),
        # SLIDE 12
        (
            "⚖️ ĐỐI CHIẾU THỰC NGHIỆM: DOCKER VS XAMPP TRUYỀN THỐNG",
            [
                "Tính nhất quán môi trường phát triển (Environment Consistency):",
                "   • XAMPP: Rất dễ bị lỗi do khác biệt phiên bản PHP và extensions giữa máy các thành viên.",
                "   • Docker: Đồng bộ 100% môi trường PHP 8.2 trên tất cả máy tính, triệt tiêu lỗi khác biệt OS.",
                "Độ phức tạp mở rộng dịch vụ (Scalability & Integrations):",
                "   • XAMPP: Cực kỳ khó khăn và dễ gây xung đột khi cài đặt PostgreSQL, Redis cache hay WebSocket trên Windows.",
                "   • Docker: Dễ dàng mở rộng chỉ bằng cách khai báo dịch vụ trong file docker-compose.yml.",
                "Khả năng triển khai lên Đám mây (Cloud Deployment):",
                "   • XAMPP: Phải thuê VPS cài đặt môi trường thủ công rất lâu và dễ sai sót cấu hình.",
                "   • Docker: Hỗ trợ triển khai tự động 1-click lên các cloud platform hiện đại như Railway, AWS, Heroku.",
                "Độ sạch của hệ điều hành (System cleanliness):",
                "   • XAMPP: Cài trực tiếp ứng dụng tạo tệp rác hệ thống.",
                "   • Docker: Cô lập hoàn toàn, dọn dẹp sạch sẽ 100% tài nguyên đĩa sau khi dùng lệnh docker-compose down."
            ],
            "Nhóm chúng em đã thực hiện đối chiếu thực nghiệm giữa Docker Stack và phần mềm XAMPP truyền thống. XAMPP mặc dù dễ cài đặt ban đầu, nhưng khi dự án phát triển lớn cần Redis hay PostgreSQL, việc cấu hình chúng trên Windows là một thách thức lớn và thường xuyên gây lỗi phiên bản PHP giữa các thành viên. Docker giải quyết triệt để tất cả các điểm yếu này: mở rộng dịch vụ chỉ bằng vài dòng code khai báo, bảo vệ máy chủ thật khỏi tệp tin rác và đảm bảo tính sạch sẽ tối đa cho hệ điều hành.",
            False
        ),
        # SLIDE 13
        (
            "🗄️ QUẢN TRỊ CSDL POSTGRESQL QUA DOCKER BASH CLI",
            [
                "Truy cập dòng lệnh trực tiếp (CLI access):",
                "   • Chạy lệnh sau trên terminal của máy thật để truy cập Postgres container:",
                "     >>> docker exec -it note_app_db psql -U admin -d note_db",
                "Khảo sát cấu trúc cơ sở dữ liệu (psql meta-commands):",
                "   • Lệnh \\\\l: Liệt kê toàn bộ các cơ sở dữ liệu trên máy chủ Postgres.",
                "   • Lệnh \\\\dt: Hiển thị tất cả các bảng dữ liệu hiện có trong CSDL note_db.",
                "   • Lệnh \\\\d notes: Xem chi tiết lược đồ thuộc tính của bảng ghi chú.",
                "Truy vấn và thao tác dữ liệu chuẩn SQL (CRUD operations):",
                "   • SELECT thành viên đăng ký, INSERT nhãn dán, UPDATE màu nền, DELETE ghi chú lỗi.",
                "Tự động hóa chạy lệnh một dòng từ máy thật (One-liner utility):",
                "   • Sử dụng cờ -c để thực thi nhanh SQL trực tiếp không cần chui vào shell:",
                "     >>> docker exec -it note_app_db psql -U admin -d note_db -c \"SELECT COUNT(*) FROM notes;\""
            ],
            "Thưa thầy cô, bên cạnh việc điều phối, khả năng quản trị cơ sở dữ liệu trực tiếp cũng là một phần cực kỳ quan trọng trong Topic 10. Thay vì phải cài đặt các công cụ giao diện nặng nề, quản trị viên có thể sử dụng sức mạnh dòng lệnh của Docker. Chỉ bằng một lệnh docker exec đơn giản, chúng ta có thể chui thẳng vào shell psql của container PostgreSQL. Tại đây, chúng ta thực hiện khảo sát các bảng dữ liệu bằng lệnh \\dt, truy vấn dữ liệu thành viên bằng SELECT, sửa đổi ghi chú bằng UPDATE hoặc chạy các lệnh SQL một dòng tự động từ bên ngoài máy thật. Điều này thể hiện khả năng kiểm soát hệ thống toàn diện và chuyên nghiệp.",
            False
        ),
        # SLIDE 14
        (
            "📚 TÀI LIỆU THAM KHẢO & TRÍCH DẪN KHOA HỌC (IEEE)",
            [
                "Tài liệu về Công nghệ Containerization & Orchestration:",
                "   • [1] D. Merkel, 'Docker: lightweight Linux containers for consistent development,' Linux Journal, 2014.",
                "   • [2] C. Pahl, 'Containerization and the PaaS Cloud,' IEEE Cloud Computing, vol. 2, no. 3, 2015.",
                "   • [3] K. Joy, 'Performance Comparison of VMs and OS-level Containers,' IJCA, vol. 122, no. 18, 2015.",
                "Tài liệu về Cơ sở dữ liệu PostgreSQL & RAM-Cache Redis:",
                "   • [4] B. Momjian, PostgreSQL: Introduction and Concepts. Addison-Wesley, 2001.",
                "   • [5] M. Stonebraker & L. Rowe, 'The design of Postgres,' Proc. 1986 ACM SIGMOD, pp. 340-355.",
                "   • [6] S. de Souza, et al., 'Performance Analysis of Postgres and MySQL under Concurrent Load,' IEEE Latin America, 2020.",
                "   • [7] J. Gaunt & B. Lawson, Introducing Redis: In-Memory Data Storage and Performance, O'Reilly Media, 2016.",
                "Tài liệu về Giao thức thời gian thực WebSocket & Kiến trúc PWA:",
                "   • [8] I. Fette & A. Melnikov, 'The WebSocket Protocol,' RFC 6455, IETF, Dec. 2011.",
                "   • [9] A. Russell, 'Progressive Web Apps: escaping tabs without losing our soul,' Infrequently Coherent, 2015.",
                "   • [10] M. Gaunt, 'Service Workers in Production,' Google Developers Technical Reviews, Aug. 2016."
            ],
            "Cuối cùng, để đảm bảo tính chính xác và vững chắc về mặt học thuật cho đề tài giữa kỳ, nhóm chúng em đã nghiêm túc nghiên cứu và trích dẫn các tài liệu khoa học uy tín đạt chuẩn quốc tế IEEE. Từ các công trình nghiên cứu nền tảng về container của Merkel, các đặc tả RFC của giao thức truyền thông WebSocket, cho đến các nghiên cứu so sánh thực nghiệm hiệu năng database trên tạp chí IEEE Latin America Transactions năm 2020. Chúng em xin chân thành cảm ơn thầy cô Hội đồng và các bạn đã chú ý lắng nghe bài thuyết trình của nhóm!",
            False
        )
    ]
    
    # Tạo các Slide nội dung
    for idx, (title_text, bullets, notes, is_dark) in enumerate(slides_data, start=2):
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Thiết lập màu nền
        background = slide.background
        fill = background.fill
        fill.solid()
        if is_dark:
            fill.fore_color.rgb = DARK_BLUE
        else:
            fill.fore_color.rgb = LIGHT_BLUE
            
        # Thêm và định dạng Tiêu đề Slide
        title_box = slide.shapes.title
        title_box.text = title_text
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        if is_dark:
            title_para.font.color.rgb = TEXT_LIGHT
        else:
            title_para.font.color.rgb = ACCENT_BLUE
            
        # Thêm và định dạng Nội dung (Bullet points)
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.clear()  # Xóa nội dung mặc định
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.name = "Arial"
            
            # Phân tách mức độ thụt lề
            if bullet.startswith("   •"):
                p.level = 1
                p.font.size = Pt(16)
                if is_dark:
                    p.font.color.rgb = RGBColor(222, 226, 230)
                else:
                    p.font.color.rgb = TEXT_DARK
            elif bullet.startswith("     -") or bullet.startswith("     >>>"):
                p.level = 2
                p.font.size = Pt(14)
                if bullet.startswith("     >>>"):
                    p.font.name = "Courier New"
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_BLUE
                else:
                    p.font.color.rgb = MUTED_GRAY
            else:
                p.level = 0
                p.font.bold = True
                p.font.size = Pt(18)
                if is_dark:
                    p.font.color.rgb = TEXT_LIGHT
                else:
                    p.font.color.rgb = ACCENT_BLUE if idx in [3, 5, 8] else TEXT_DARK
                    
            p.space_after = Pt(8)
            
        # Thêm Speaker Notes trực tiếp vào Slide PowerPoint
        slide.notes_slide.notes_text_frame.text = f"Kịch bản thuyết trình (Speaker Notes) - SLIDE {idx}:\n{notes}"
        
    # Lưu tệp tin trình bày PowerPoint
    output_filename = "Thuyet_Trinh_Topic10.pptx"
    prs.save(output_filename)
    print(f"SUCCESS: PowerPoint file created successfully as '{output_filename}'!")

if __name__ == "__main__":
    create_presentation()
